import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
LOGO_PATH  = os.path.join(SCRIPT_DIR, "mycom-logo.png")

# ── Colour palette ──────────────────────────────────────────────
DARK_BG      = RGBColor(0x1B, 0x1F, 0x3B)
ACCENT_BLUE  = RGBColor(0x00, 0x9E, 0xFF)
ACCENT_GREEN = RGBColor(0x00, 0xE6, 0x76)
ACCENT_RED   = RGBColor(0xFF, 0x4D, 0x4D)
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY   = RGBColor(0xCC, 0xCC, 0xCC)
MID_GRAY     = RGBColor(0x99, 0x99, 0x99)
SECTION_BG   = RGBColor(0x23, 0x27, 0x4A)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SLIDE_W = prs.slide_width
SLIDE_H = prs.slide_height

# ── helpers ─────────────────────────────────────────────────────
def solid_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Segoe UI"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_bullets(slide, left, top, width, height, items, font_size=16,
                color=WHITE, bullet_color=ACCENT_BLUE, spacing=Pt(8), font_name="Segoe UI"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = font_name
        p.space_after = spacing
        p.level = 0
    return txBox

def slide_number_footer(slide, num):
    add_text_box(slide, SLIDE_W - Inches(1), SLIDE_H - Inches(0.5),
                 Inches(0.8), Inches(0.4), str(num),
                 font_size=11, color=MID_GRAY, alignment=PP_ALIGN.RIGHT)

def accent_bar(slide, top=Inches(1.35)):
    add_rect(slide, Inches(0.8), top, Inches(2), Pt(4), ACCENT_BLUE)

def add_mycom_footer(slide):
    """Add MYCOM logo + '© MYCOM Confidential' footer to a slide."""
    slide.shapes.add_picture(LOGO_PATH,
                             Inches(0.85), Inches(7.03),
                             Inches(1.84), Inches(0.30))
    add_text_box(slide, Inches(0.67), Inches(6.97), Inches(12), Inches(0.33),
                 "\u00a9 MYCOM Confidential",
                 font_size=13, color=MID_GRAY, alignment=PP_ALIGN.CENTER,
                 font_name="Segoe UI")

def add_mycom_title_branding(slide):
    """Add large MYCOM logo + confidential disclaimer to the title slide."""
    slide.shapes.add_picture(LOGO_PATH,
                             Inches(0.83), Inches(0.56),
                             Inches(3.57), Inches(0.59))
    add_text_box(slide, Inches(0.67), Inches(6.3), Inches(12), Inches(0.6),
                 "CONFIDENTIAL INFORMATION, RESTRICTED TO MYCOM",
                 font_size=10, color=MID_GRAY, alignment=PP_ALIGN.CENTER,
                 font_name="Segoe UI")

def section_title(slide, title, subtitle=None, slide_num=1):
    solid_bg(slide, DARK_BG)
    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.7),
                 title, font_size=32, color=WHITE, bold=True)
    accent_bar(slide)
    if subtitle:
        add_text_box(slide, Inches(0.8), Inches(1.55), Inches(11), Inches(0.5),
                     subtitle, font_size=16, color=LIGHT_GRAY)
    slide_number_footer(slide, slide_num)
    add_mycom_footer(slide)

# ── SLIDE 1 — Title ────────────────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6])  # blank
solid_bg(s, DARK_BG)
add_rect(s, 0, 0, SLIDE_W, Inches(0.12), ACCENT_BLUE)
add_text_box(s, Inches(1.5), Inches(1.8), Inches(10), Inches(1.2),
             "Non-Blocking Bulk Operations\nManagement using SSE",
             font_size=40, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_rect(s, Inches(5.5), Inches(3.3), Inches(2.3), Pt(3), ACCENT_BLUE)
add_text_box(s, Inches(1.5), Inches(3.6), Inches(10), Inches(0.8),
             "Real-time progress tracking for long-running telecom operations\nwith Server-Sent Events, React & Spring Boot",
             font_size=18, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
add_text_box(s, Inches(1.5), Inches(5.5), Inches(10), Inches(0.5),
             "Architecture & Design Walkthrough",
             font_size=14, color=MID_GRAY, alignment=PP_ALIGN.CENTER)
slide_number_footer(s, 1)
add_mycom_title_branding(s)
add_mycom_footer(s)

# ── SLIDE 2 — Problem Statement ────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6])
section_title(s, "Problem Statement", "Why traditional synchronous APIs fall short", 2)

problems = [
    ("UI Blocking", "Synchronous REST calls lock the browser thread during\nlong-running bulk operations (network create/update/delete)."),
    ("Poor User Experience", "Users see a spinner for 1-3 minutes with no progress\ninformation, leading to repeated clicks or tab closures."),
    ("Timeout Risk", "HTTP connections may timeout before the operation completes,\nleaving the user unaware of the actual result."),
    ("No Progress Visibility", "Traditional request-response gives zero intermediate\nfeedback — it is all-or-nothing."),
]
y = Inches(2.2)
for title, desc in problems:
    add_rect(s, Inches(0.8), y, Inches(11.5), Inches(1.05), SECTION_BG, ACCENT_BLUE)
    add_text_box(s, Inches(1.1), y + Pt(6), Inches(3), Inches(0.4),
                 title, font_size=18, color=ACCENT_BLUE, bold=True)
    add_text_box(s, Inches(4.0), y + Pt(4), Inches(8), Inches(0.9),
                 desc, font_size=14, color=LIGHT_GRAY)
    y += Inches(1.2)

# ── SLIDE 3 — Traditional Approaches ──────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6])
section_title(s, "Traditional Approaches & Limitations", None, 3)

approaches = [
    ("Blocking REST", [
        "Client waits for full response — 1-3 min freeze",
        "Single timeout kills entire operation",
        "No intermediate progress feedback",
    ], ACCENT_RED),
    ("Polling", [
        "Client repeatedly hits GET /status every N seconds",
        "Wastes bandwidth & server resources",
        "Delayed feedback (up to N-second lag)",
    ], RGBColor(0xFF, 0xAA, 0x00)),
    ("WebSockets", [
        "Full-duplex — overkill for unidirectional server push",
        "Complex connection lifecycle management",
        "Harder to scale behind load balancers & proxies",
    ], RGBColor(0xFF, 0xAA, 0x00)),
]
x = Inches(0.8)
for title, items, clr in approaches:
    add_rect(s, x, Inches(2.0), Inches(3.7), Inches(4.5), SECTION_BG, clr)
    add_text_box(s, x + Inches(0.3), Inches(2.15), Inches(3.2), Inches(0.5),
                 title, font_size=20, color=clr, bold=True)
    add_bullets(s, x + Inches(0.3), Inches(2.8), Inches(3.2), Inches(3.5),
                [f"• {i}" for i in items], font_size=13, color=LIGHT_GRAY)
    x += Inches(4.0)

add_text_box(s, Inches(0.8), Inches(6.6), Inches(11), Inches(0.5),
             "Verdict: SSE provides lightweight, server-to-client streaming with automatic reconnection — ideal for progress updates.",
             font_size=14, color=ACCENT_GREEN, bold=True)

# ── SLIDE 4 — Solution Overview ────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6])
section_title(s, "Solution Overview", "Non-blocking architecture with SSE-based progress streaming", 4)

pillars = [
    ("Fire-and-Forget", "Client sends command,\ngets 202 Accepted\nwith a Task ID instantly."),
    ("Async Execution", "Spring Boot processes\nthe operation on a\ndedicated thread pool."),
    ("SSE Streaming", "Real-time progress\nevents pushed to the\nclient via EventSource."),
    ("Optimistic UI", "React updates the UI\nimmediately — pending\nstate shown until done."),
]
x = Inches(0.6)
for title, desc in pillars:
    add_rect(s, x, Inches(2.3), Inches(2.9), Inches(3.0), SECTION_BG, ACCENT_BLUE)
    add_text_box(s, x + Inches(0.2), Inches(2.5), Inches(2.5), Inches(0.5),
                 title, font_size=18, color=ACCENT_BLUE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(s, x + Inches(0.2), Inches(3.2), Inches(2.5), Inches(1.8),
                 desc, font_size=14, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
    x += Inches(3.15)

add_text_box(s, Inches(0.8), Inches(5.8), Inches(11), Inches(1),
             "Key Result: The user interface remains fully responsive throughout multi-minute operations,\n"
             "with a live progress bar and contextual messages at every stage.",
             font_size=14, color=WHITE)

# ── SLIDE 5 — High-Level Architecture ─────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6])
section_title(s, "High-Level Architecture", "End-to-end component map", 5)

# Component boxes
components = [
    (Inches(0.5),  Inches(2.8), Inches(3), Inches(2.2), "React Frontend",
     "• NetworkList, NetworkForm\n• TaskContext (global state)\n• SSEManager (EventSource)\n• Mock-data fallback", ACCENT_BLUE),
    (Inches(5.1),  Inches(2.3), Inches(3.2), Inches(1.3), "Spring Boot — REST",
     "• NetworkController (CRUD)\n• TaskController (SSE stream)", ACCENT_GREEN),
    (Inches(5.1),  Inches(4.0), Inches(3.2), Inches(1.3), "Spring Boot — Async",
     "• TaskService (@Async)\n• ThreadPool: 5-10 threads", RGBColor(0xFF, 0xAA, 0x00)),
    (Inches(9.8),  Inches(2.8), Inches(3), Inches(2.2), "MySQL Database",
     "• networks table\n• JPA / Hibernate ORM\n• Status enum constraint\n• ddl-auto=validate", ACCENT_RED),
]
for lf, tp, w, h, title, desc, clr in components:
    add_rect(s, lf, tp, w, h, SECTION_BG, clr)
    add_text_box(s, lf + Inches(0.15), tp + Pt(4), w - Inches(0.3), Inches(0.4),
                 title, font_size=16, color=clr, bold=True)
    add_text_box(s, lf + Inches(0.15), tp + Inches(0.45), w - Inches(0.3), h - Inches(0.5),
                 desc, font_size=12, color=LIGHT_GRAY)

# Arrows as text labels
add_text_box(s, Inches(3.5), Inches(2.5), Inches(1.6), Inches(0.4),
             "REST + SSE  -->", font_size=12, color=ACCENT_BLUE, bold=True)
add_text_box(s, Inches(8.3), Inches(3.2), Inches(1.6), Inches(0.4),
             "JPA  -->", font_size=12, color=ACCENT_GREEN, bold=True)

add_text_box(s, Inches(0.5), Inches(5.8), Inches(12), Inches(1),
             "Communication: REST (commands) + SSE (progress) | CORS enabled for localhost:3000 --> localhost:8080",
             font_size=13, color=MID_GRAY)

# ── SLIDE 6 — Backend Design ──────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6])
section_title(s, "Backend Design — Spring Boot", "Controllers, async services, SSE emitter management", 6)

left_items = [
    "NetworkController — synchronous reads (GET), async writes (POST/PUT/DELETE return 202)",
    "TaskController — exposes GET /api/tasks/{taskId}/progress as text/event-stream",
    "NetworkService — CRUD via Spring Data JPA (NetworkRepository)",
    "TaskService — creates tasks, manages ConcurrentHashMap<taskId, SseEmitter>",
]
right_items = [
    "@Async(\"taskExecutor\") — methods run on dedicated thread pool",
    "ThreadPoolTaskExecutor: core=5, max=10, queue=25, prefix=\"NetworkTask-\"",
    "SseEmitter timeout: 300 000 ms (5 minutes)",
    "Lifecycle callbacks: onCompletion / onTimeout / onError clean up emitter map",
    "Events emitted: default (progress), 'complete', 'error'",
]
add_text_box(s, Inches(0.8), Inches(1.8), Inches(5.5), Inches(0.4),
             "Controllers & Services", font_size=18, color=ACCENT_BLUE, bold=True)
add_bullets(s, Inches(0.8), Inches(2.3), Inches(5.5), Inches(4.5),
            [f"• {i}" for i in left_items], font_size=13, color=LIGHT_GRAY)

add_text_box(s, Inches(6.8), Inches(1.8), Inches(5.5), Inches(0.4),
             "Threading & SSE Model", font_size=18, color=ACCENT_GREEN, bold=True)
add_bullets(s, Inches(6.8), Inches(2.3), Inches(5.8), Inches(4.5),
            [f"• {i}" for i in right_items], font_size=13, color=LIGHT_GRAY)

# ── SLIDE 7 — Database Design ─────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6])
section_title(s, "Database Design — MySQL", "Schema, state persistence, and ORM mapping", 7)

cols = [
    "id  VARCHAR(50)  PK — format: net-{UUID-8}",
    "name  VARCHAR(255)  NOT NULL",
    "type  VARCHAR(50)  — 5G, 4G LTE, NB-IoT, etc.",
    "status  VARCHAR(11)  — ACTIVE | STANDBY | MAINTENANCE | INACTIVE",
    "region  VARCHAR(50)",
    "ip_range  VARCHAR(20) — CIDR notation",
    "bandwidth  VARCHAR(20)  |  latency  VARCHAR(10)  |  nodes  INT",
    "created_at  TIMESTAMP  |  updated_at  TIMESTAMP",
    "description  TEXT",
]
add_text_box(s, Inches(0.8), Inches(1.8), Inches(5), Inches(0.4),
             "Table: networks", font_size=20, color=ACCENT_BLUE, bold=True)
add_bullets(s, Inches(0.8), Inches(2.3), Inches(6), Inches(4.5),
            [f"• {c}" for c in cols], font_size=13, color=LIGHT_GRAY, spacing=Pt(5))

why_items = [
    "DB-backed state survives backend restarts (unlike in-memory task state)",
    "Hibernate validates schema on startup (ddl-auto=validate)",
    "NetworkRepository extends JpaRepository — findByStatus(enum) derived query",
    "Note: TaskInfo is in-memory only (ConcurrentHashMap) — lost on restart",
    "Future: persist tasks to DB for resume-after-crash capability",
]
add_text_box(s, Inches(7.3), Inches(1.8), Inches(5), Inches(0.4),
             "Why DB-Backed State?", font_size=20, color=ACCENT_GREEN, bold=True)
add_bullets(s, Inches(7.3), Inches(2.3), Inches(5.5), Inches(4.5),
            [f"• {w}" for w in why_items], font_size=13, color=LIGHT_GRAY, spacing=Pt(5))

# ── SLIDE 8 — SSE Workflow ─────────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6])
section_title(s, "SSE Workflow — Sequence Flow", "From user click to completion event", 8)

# Flowchart: snake pattern — 2 rows × 4 shapes
# Row 1 (left → right):  1  →  2  →  3  →  4
#                                              ↓
# Row 2 (right → left):  8  ←  7  ←  6  ←  5

FC_W   = Inches(2.3)
FC_H   = Inches(1.2)
FC_COL = [Inches(0.6), Inches(3.5), Inches(6.4), Inches(9.3)]
FC_R1  = Inches(2.0)
FC_R2  = Inches(4.7)

flow_steps = [
    # (num, title, description, flowchart_shape, accent_color)
    ("1", "User Action",      "Click Create / Update\n/ Delete in React UI",   MSO_SHAPE.FLOWCHART_TERMINATOR,       ACCENT_BLUE),
    ("2", "REST Call (202)",   "POST /api/networks\n→ 202 Accepted + taskId",  MSO_SHAPE.FLOWCHART_DATA,             ACCENT_BLUE),
    ("3", "SSE Subscribe",    "Open EventSource to\n/api/tasks/{id}/progress", MSO_SHAPE.FLOWCHART_PROCESS,          ACCENT_BLUE),
    ("4", "Async Execution",  "@Async thread starts\n1-3 min duration",        MSO_SHAPE.FLOWCHART_PREDEFINED_PROCESS, ACCENT_BLUE),
    ("5", "Progress Events",  "SseEmitter.send()\nevery 3 seconds",            MSO_SHAPE.FLOWCHART_MULTIDOCUMENT,    ACCENT_BLUE),
    ("6", "DB Persist",       "networkRepository\n.save() → MySQL",            MSO_SHAPE.FLOWCHART_MAGNETIC_DISK,    ACCENT_GREEN),
    ("7", "Complete Event",   "SSE 'complete' event\nstream closes",            MSO_SHAPE.FLOWCHART_DOCUMENT,         ACCENT_GREEN),
    ("8", "UI Refresh",       "TASK_COMPLETED\nfetch fresh data",              MSO_SHAPE.FLOWCHART_TERMINATOR,       ACCENT_GREEN),
]

def _add_fc_shape(slide, shp_type, left, top, w, h, clr, num, title, desc):
    shape = slide.shapes.add_shape(shp_type, left, top, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = SECTION_BG
    shape.line.color.rgb = clr
    shape.line.width = Pt(2)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = f"{num}. {title}"
    p.font.size = Pt(12)
    p.font.color.rgb = clr
    p.font.bold = True
    p.font.name = "Segoe UI"
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = desc
    p2.font.size = Pt(9)
    p2.font.color.rgb = LIGHT_GRAY
    p2.font.name = "Segoe UI"
    p2.alignment = PP_ALIGN.CENTER

def _add_arrow(slide, shape_type, left, top, w, h, clr):
    arrow = slide.shapes.add_shape(shape_type, left, top, w, h)
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = clr
    arrow.line.fill.background()

# ── Row 1: steps 1-4 (left → right) ──
for i in range(4):
    num, title, desc, shp_type, clr = flow_steps[i]
    _add_fc_shape(s, shp_type, FC_COL[i], FC_R1, FC_W, FC_H, clr, num, title, desc)

# Right arrows between row-1 shapes
for i in range(3):
    ax = FC_COL[i] + FC_W
    gap = FC_COL[i + 1] - FC_COL[i] - FC_W
    _add_arrow(s, MSO_SHAPE.RIGHT_ARROW,
               ax, FC_R1 + FC_H // 2 - Inches(0.13), gap, Inches(0.26), ACCENT_BLUE)

# ── Down arrow: step 4 → step 5 ──
_add_arrow(s, MSO_SHAPE.DOWN_ARROW,
           FC_COL[3] + FC_W // 2 - Inches(0.13),
           FC_R1 + FC_H,
           Inches(0.26), FC_R2 - FC_R1 - FC_H,
           ACCENT_BLUE)

# ── Row 2: steps 5-8 (right → left) ──
for i in range(4):
    num, title, desc, shp_type, clr = flow_steps[4 + i]
    col_idx = 3 - i
    _add_fc_shape(s, shp_type, FC_COL[col_idx], FC_R2, FC_W, FC_H, clr, num, title, desc)

# Left arrows between row-2 shapes
for i in range(3):
    src_col = 3 - i
    dst_col = src_col - 1
    ax = FC_COL[dst_col] + FC_W
    gap = FC_COL[src_col] - FC_COL[dst_col] - FC_W
    clr = ACCENT_GREEN if i >= 1 else ACCENT_BLUE
    _add_arrow(s, MSO_SHAPE.LEFT_ARROW,
               ax, FC_R2 + FC_H // 2 - Inches(0.13), gap, Inches(0.26), clr)

# ── SLIDE 9 — Frontend Design ─────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6])
section_title(s, "Frontend Design — React", "Component architecture, SSE integration, mock fallback", 9)

left = [
    "NetworkList — grid of network cards with pending-state overlay",
    "NetworkForm — validated create/edit modal (CIDR regex, required fields)",
    "NetworkDetails — full view with edit/delete actions",
    "TaskProgressPanel — floating panel showing all active + recent tasks",
    "TaskItem + ProgressBar — per-task live progress display",
]
right = [
    "SSEManager class wraps EventSource with reconnect logic",
    "Exponential backoff: 1s -> 2s -> 4s -> 8s -> 16s (max 5 retries)",
    "TaskContext (React Context + useReducer) — global task state",
    "Actions: TASK_STARTED, TASK_PROGRESS, TASK_COMPLETED, TASK_FAILED",
    "networkService.js — tries backend, falls back to mockNetworks.js",
    "REACT_APP_USE_MOCK=true activates MockSSEConnection (10% failure sim)",
]
add_text_box(s, Inches(0.8), Inches(1.8), Inches(5), Inches(0.4),
             "UI Components", font_size=18, color=ACCENT_BLUE, bold=True)
add_bullets(s, Inches(0.8), Inches(2.3), Inches(5.5), Inches(4.5),
            [f"• {i}" for i in left], font_size=13, color=LIGHT_GRAY, spacing=Pt(5))

add_text_box(s, Inches(6.8), Inches(1.8), Inches(5.5), Inches(0.4),
             "SSE & State Management", font_size=18, color=ACCENT_GREEN, bold=True)
add_bullets(s, Inches(6.8), Inches(2.3), Inches(5.8), Inches(4.5),
            [f"• {i}" for i in right], font_size=13, color=LIGHT_GRAY, spacing=Pt(5))

# ── SLIDE 10 — Error Handling & Resilience ─────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6])
section_title(s, "Error Handling & Resilience", "Graceful degradation at every layer", 10)

rows = [
    ("Backend Failure",     "Async try-catch wraps every operation\nfailTask() emits SSE error event with details\nTask state preserved for audit / retry"),
    ("SSE Disconnect",      "Exponential backoff reconnect (5 retries)\nonMaxRetriesExceeded callback notifies UI\n5-minute server-side emitter timeout"),
    ("Backend Unavailable", "networkService.js catches fetch errors + 5s AbortController timeout\nAutomatically returns mock data with source='MOCK'\nUI shows 'Offline — showing cached data' badge"),
    ("Optimistic Rollback", "Create/Update shown immediately with _pending flag\nOn failure: pending overlay stays, user can retry\nOn success: flag removed, fresh data fetched from DB"),
]
y = Inches(2.0)
for title, desc in rows:
    add_rect(s, Inches(0.8), y, Inches(11.5), Inches(1.15), SECTION_BG, ACCENT_RED)
    add_text_box(s, Inches(1.1), y + Pt(6), Inches(3), Inches(0.4),
                 title, font_size=16, color=ACCENT_RED, bold=True)
    add_text_box(s, Inches(4.2), y + Pt(2), Inches(7.8), Inches(1.05),
                 desc, font_size=13, color=LIGHT_GRAY)
    y += Inches(1.3)

# ── SLIDE 11 — Limitations of SSE ─────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6])
section_title(s, "Limitations of SSE", "Trade-offs to consider when choosing Server-Sent Events", 11)

limitations = [
    ("Unidirectional Only",
     "SSE is server-to-client only — the client cannot send data\n"
     "back over the same connection. A separate REST call is needed\n"
     "for every client action."),
    ("Browser Connection Limit",
     "HTTP/1.1 allows only ~6 concurrent connections per domain.\n"
     "Each SSE stream occupies one, starving other requests.\n"
     "HTTP/2 multiplexing mitigates this."),
    ("Text Only — No Binary",
     "SSE transmits UTF-8 text exclusively. Binary payloads\n"
     "(images, protobuf) must be Base64-encoded, adding ~33% overhead."),
    ("No Custom Headers",
     "The browser EventSource API does not support custom HTTP\n"
     "headers (e.g., Authorization: Bearer). Tokens must be passed\n"
     "via query params or cookies."),
    ("Proxy & Firewall Buffering",
     "Corporate proxies and CDNs may buffer or terminate long-lived\n"
     "HTTP connections, breaking the event stream silently."),
]
y = Inches(2.1)
for title, desc in limitations:
    add_rect(s, Inches(0.8), y, Inches(11.5), Inches(0.88), SECTION_BG, RGBColor(0xFF, 0xAA, 0x00))
    add_text_box(s, Inches(1.1), y + Pt(6), Inches(3.2), Inches(0.4),
                 title, font_size=16, color=RGBColor(0xFF, 0xAA, 0x00), bold=True)
    add_text_box(s, Inches(4.5), y + Pt(2), Inches(7.5), Inches(0.8),
                 desc, font_size=13, color=LIGHT_GRAY)
    y += Inches(0.95)

# ── SLIDE 12 — Benefits & Outcomes ─────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6])
section_title(s, "Benefits & Outcomes", None, 12)

benefits = [
    ("Responsive UX",
     "Zero UI blocking — users interact freely during\nmulti-minute operations with live progress feedback."),
    ("Scalable Async Processing",
     "ThreadPoolTaskExecutor isolates long-running work.\n5 core + 10 max threads with 25-task queue."),
    ("Clean Separation of Concerns",
     "REST for commands, SSE for events, React Context\nfor state — each layer independently testable."),
    ("Resilient by Default",
     "Automatic mock fallback, SSE reconnect, optimistic\nUI — the app degrades gracefully, never crashes."),
    ("Production-Ready Patterns",
     "ConcurrentHashMap, @Async, SseEmitter lifecycle\ncallbacks, CORS config, enum-safe status handling."),
]
y = Inches(2.0)
for title, desc in benefits:
    add_rect(s, Inches(0.8), y, Inches(0.12), Inches(0.85), ACCENT_GREEN)
    add_text_box(s, Inches(1.2), y + Pt(2), Inches(3.5), Inches(0.4),
                 title, font_size=17, color=ACCENT_GREEN, bold=True)
    add_text_box(s, Inches(5.0), y + Pt(2), Inches(7.5), Inches(0.8),
                 desc, font_size=14, color=LIGHT_GRAY)
    y += Inches(0.95)

# ── SLIDE 13 — Future Enhancements ────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6])
section_title(s, "Future Enhancements", "Roadmap for production hardening", 13)

enhancements = [
    ("Retry & Resume",         "Persist TaskInfo to DB so operations can resume after crash or restart"),
    ("Progress Across Refresh", "Store task state in localStorage / sessionStorage for cross-tab continuity"),
    ("Auth & Security",        "JWT/OAuth2 on REST + SSE endpoints; role-based operation permissions"),
    ("Message Broker",         "Kafka or RabbitMQ decouples task creation from execution — enables horizontal scaling"),
    ("Bulk Selection",         "Multi-select networks for batch delete/update as a single tracked operation"),
    ("Observability",          "Structured logging, Micrometer metrics, distributed tracing (OpenTelemetry)"),
]
y = Inches(2.0)
for title, desc in enhancements:
    add_rect(s, Inches(0.8), y, Inches(11.5), Inches(0.72), SECTION_BG)
    add_text_box(s, Inches(1.1), y + Pt(4), Inches(3.2), Inches(0.4),
                 title, font_size=16, color=ACCENT_BLUE, bold=True)
    add_text_box(s, Inches(4.5), y + Pt(4), Inches(7.5), Inches(0.5),
                 desc, font_size=14, color=LIGHT_GRAY)
    y += Inches(0.82)

# ── Save ────────────────────────────────────────────────────────
out_path = r"C:\Users\hdahiya1\Downloads\Mycom Project\Non_Blocking_Operations_SSE_Presentation.pptx"
prs.save(out_path)
print(f"Saved to: {out_path}")
